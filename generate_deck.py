import os
import shutil
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# Paths for generated mockup files
DASHBOARD_SRC = r"C:\Users\DELL\.gemini\antigravity-ide\brain\2dcd6c07-98d1-4f04-a54f-e5e271253122\mockup_dashboard_1781775207502.png"
ACCESSIBILITY_SRC = r"C:\Users\DELL\.gemini\antigravity-ide\brain\2dcd6c07-98d1-4f04-a54f-e5e271253122\mockup_accessibility_1781775223261.png"

# Target local paths
DASHBOARD_DEST = "mockup_dashboard.png"
ACCESSIBILITY_DEST = "mockup_accessibility.png"

# Copy images to local directory for ease of reference
try:
    if os.path.exists(DASHBOARD_SRC):
        shutil.copy(DASHBOARD_SRC, DASHBOARD_DEST)
        print(f"Copied dashboard mockup to {DASHBOARD_DEST}")
    else:
        print(f"Warning: Source dashboard mockup not found at {DASHBOARD_SRC}")
        
    if os.path.exists(ACCESSIBILITY_SRC):
        shutil.copy(ACCESSIBILITY_SRC, ACCESSIBILITY_DEST)
        print(f"Copied accessibility mockup to {ACCESSIBILITY_DEST}")
    else:
        print(f"Warning: Source accessibility mockup not found at {ACCESSIBILITY_SRC}")
except Exception as e:
    print(f"Error copying mockup assets: {e}")

# Color Scheme
DARK_BLUE = RGBColor(10, 25, 47)      # Premium background/theme color #0a192f
DARK_BLUE_LIGHT = RGBColor(23, 42, 69) # Card background color for dark mode #172a45
EMERALD = RGBColor(16, 185, 129)      # Primary accent color #10b981
EMERALD_MUTED = RGBColor(5, 150, 105)  # Secondary accent color #059669
WHITE = RGBColor(255, 255, 255)
OFF_WHITE = RGBColor(248, 250, 252)   # Standard background #f8fafc
TEXT_DARK = RGBColor(30, 41, 59)      # Slate dark text #1e293b
TEXT_MUTED = RGBColor(100, 116, 139)  # Slate secondary text #64748b

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6] # Blank slide

def add_full_background(slide, color):
    """Draws a full slide background rectangle to replace default white."""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.color.rgb = color
    return bg

def add_header(slide, title_text, dark_mode=False):
    """Adds a consistent header banner to the slide."""
    # Top Accent Bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = EMERALD
    bar.line.color.rgb = EMERALD
    
    # Title Text Box
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.0), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = "Arial"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE if dark_mode else DARK_BLUE

def add_card(slide, left, top, width, height, bg_color, border_color=None):
    """Creates a stylized background card with custom fills and borders."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    if border_color:
        card.line.color.rgb = border_color
        card.line.width = Pt(1.5)
    else:
        card.line.color.rgb = bg_color
    return card

# ==========================================
# SLIDE 1: TITLE SLIDE (DARK MODE)
# ==========================================
slide1 = prs.slides.add_slide(blank_layout)
add_full_background(slide1, DARK_BLUE)

# Bottom Emerald line
line = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.3), Inches(13.333), Inches(0.2))
line.fill.solid()
line.fill.fore_color.rgb = EMERALD
line.line.color.rgb = EMERALD

# Title and Subtitle Text Frame
title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.3), Inches(2.5))
tf = title_box.text_frame
tf.word_wrap = True

# AcademixIQ Title
p1 = tf.paragraphs[0]
p1.text = "AcademixIQ"
p1.font.name = "Arial"
p1.font.size = Pt(64)
p1.font.bold = True
p1.font.color.rgb = WHITE
p1.space_after = Pt(10)

# Subtitle
p2 = tf.add_paragraph()
p2.text = "AI-Powered Smart Learning Companion & Accessibility Hub"
p2.font.name = "Arial"
p2.font.size = Pt(22)
p2.font.color.rgb = EMERALD
p2.space_after = Pt(20)

# Hackathon tag
p3 = tf.add_paragraph()
p3.text = "Bharat Academix CodeQuest 2026  |  Round 1: Idea Proposal"
p3.font.name = "Arial"
p3.font.size = Pt(16)
p3.font.color.rgb = WHITE

# Team Info Card (Bottom Left)
team_box = slide1.shapes.add_textbox(Inches(1.0), Inches(5.2), Inches(6.0), Inches(1.5))
tf_team = team_box.text_frame
tf_team.word_wrap = True
pt1 = tf_team.paragraphs[0]
pt1.text = "Team: ayushkhaitan2004"
pt1.font.name = "Arial"
pt1.font.bold = True
pt1.font.size = Pt(14)
pt1.font.color.rgb = EMERALD
pt2 = tf_team.add_paragraph()
pt2.text = "Solo Innovator: Ayush Khaitan  |  SRMIST Ghaziabad"
pt2.font.name = "Arial"
pt2.font.size = Pt(13)
pt2.font.color.rgb = WHITE
pt3 = tf_team.add_paragraph()
pt3.text = "Themes: EdTech, AI/ML, Accessibility & Inclusion"
pt3.font.name = "Arial"
pt3.font.size = Pt(11)
pt3.font.color.rgb = TEXT_MUTED

# ==========================================
# SLIDE 2: THE PROBLEM STATEMENT (LIGHT MODE)
# ==========================================
slide2 = prs.slides.add_slide(blank_layout)
add_full_background(slide2, OFF_WHITE)
add_header(slide2, "The Problem Statement")

# Left Column - The Challenges Card
add_card(slide2, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.2), WHITE, border_color=EMERALD)
challenges_box = slide2.shapes.add_textbox(Inches(1.1), Inches(1.7), Inches(5.0), Inches(4.8))
tf_c = challenges_box.text_frame
tf_c.word_wrap = True

pc1 = tf_c.paragraphs[0]
pc1.text = "Key Challenges in Indian Education:"
pc1.font.name = "Arial"
pc1.font.size = Pt(20)
pc1.font.bold = True
pc1.font.color.rgb = DARK_BLUE
pc1.space_after = Pt(18)

points_c = [
    ("Passive Learning Model:", " Traditional systems rely on one-way lectures, leading to poor retention and zero real-time personalized feedback."),
    ("Regional Language Barriers:", " 80%+ of higher academic research, papers, and specialized resources are exclusively in English, isolating regional-medium students."),
    ("Accessibility Gaps:", " Neurodivergent learners (ADHD, Dyslexia) and visually/hearing-impaired students lack interactive, adaptive support in physical/digital classrooms.")
]

for title, desc in points_c:
    p = tf_c.add_paragraph()
    p.space_after = Pt(12)
    run_title = p.add_run()
    run_title.text = "• " + title
    run_title.font.bold = True
    run_title.font.size = Pt(13)
    run_title.font.color.rgb = DARK_BLUE
    run_desc = p.add_run()
    run_desc.text = desc
    run_desc.font.size = Pt(13)
    run_desc.font.color.rgb = TEXT_DARK

# Right Column - Data & Stat Card
add_card(slide2, Inches(6.8), Inches(1.5), Inches(5.6), Inches(5.2), DARK_BLUE)
stat_box = slide2.shapes.add_textbox(Inches(7.1), Inches(1.7), Inches(5.0), Inches(4.8))
tf_s = stat_box.text_frame
tf_s.word_wrap = True

ps1 = tf_s.paragraphs[0]
ps1.text = "The Opportunity in Numbers:"
ps1.font.name = "Arial"
ps1.font.size = Pt(20)
ps1.font.bold = True
ps1.font.color.rgb = EMERALD
ps1.space_after = Pt(24)

stats = [
    ("250 Million+", "Students enrolled in Indian primary/secondary schools needing tailored education."),
    ("< 5%", "Of educational content conforms to international accessibility standards in regional domains."),
    ("80%", "Of classroom knowledge loss occurs within 48 hours without active reinforcement.")
]

for stat, desc in stats:
    p = tf_s.add_paragraph()
    p.space_after = Pt(14)
    run_stat = p.add_run()
    run_stat.text = stat + "\n"
    run_stat.font.bold = True
    run_stat.font.size = Pt(24)
    run_stat.font.color.rgb = WHITE
    run_desc = p.add_run()
    run_desc.text = desc
    run_desc.font.size = Pt(12)
    run_desc.font.color.rgb = EMERALD

# ==========================================
# SLIDE 3: THE SOLUTION - ACADEMIXIQ (LIGHT MODE)
# ==========================================
slide3 = prs.slides.add_slide(blank_layout)
add_full_background(slide3, OFF_WHITE)
add_header(slide3, "The Proposed Solution - AcademixIQ")

# Theme Alignment Footer Banner
theme_lbl = slide3.shapes.add_textbox(Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.5))
tf_theme = theme_lbl.text_frame
tf_theme.word_wrap = True
pth = tf_theme.paragraphs[0]
pth.text = "Direct Alignment: Education Technology  |  Social Impact & Accessibility  |  AI & Machine Learning"
pth.font.bold = True
pth.font.size = Pt(11)
pth.font.color.rgb = EMERALD_MUTED
pth.alignment = PP_ALIGN.CENTER

# 4 pillar card layouts
pillars = [
    ("BhashaTranslate", "Real-time vernacular lecture transcription, translating English classes to local Indian languages (Hindi, Tamil, etc.) instantly.", Inches(0.8), Inches(1.6)),
    ("VisualMind Canvas", "Generates real-time, interactive mind maps and visual flowcharts from lecture audio or PDFs to simplify complex concepts.", Inches(6.8), Inches(1.6)),
    ("SkillBridge AI", "Contextual, adaptive quiz generation that analyzes user performance to pinpoint concept gaps and customize roadmaps.", Inches(0.8), Inches(4.2)),
    ("EduAccess Engine", "High-contrast screen-reader styles, speech-to-text navigation, and simplified content modules for neurodiverse learners.", Inches(6.8), Inches(4.2))
]

for title, desc, left, top in pillars:
    add_card(slide3, left, top, Inches(5.6), Inches(2.2), WHITE, border_color=EMERALD)
    tb = slide3.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), Inches(5.2), Inches(1.8))
    tf_p = tb.text_frame
    tf_p.word_wrap = True
    
    p1 = tf_p.paragraphs[0]
    p1.text = title
    p1.font.name = "Arial"
    p1.font.bold = True
    p1.font.size = Pt(18)
    p1.font.color.rgb = DARK_BLUE
    p1.space_after = Pt(8)
    
    p2 = tf_p.add_paragraph()
    p2.text = desc
    p2.font.name = "Arial"
    p2.font.size = Pt(13)
    p2.font.color.rgb = TEXT_DARK

# ==========================================
# SLIDE 4: DASHBOARD SHOWCASE (DARK MODE)
# ==========================================
slide4 = prs.slides.add_slide(blank_layout)
add_full_background(slide4, DARK_BLUE)
add_header(slide4, "Interactive Student Dashboard Mockup", dark_mode=True)

# Left text box
desc_box = slide4.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(4.8), Inches(5.0))
tf_d = desc_box.text_frame
tf_d.word_wrap = True

pd1 = tf_d.paragraphs[0]
pd1.text = "AcademixIQ Student Dashboard:"
pd1.font.name = "Arial"
pd1.font.size = Pt(22)
pd1.font.bold = True
pd1.font.color.rgb = EMERALD
pd1.space_after = Pt(16)

features_d = [
    ("Glassmorphism UI:", " Optimized for focus with a modern dark aesthetics design system."),
    ("Visual Mind-Mapping:", " Center stage is occupied by interactive node diagrams illustrating the relation between academic topics."),
    ("Real-Time Transcripts:", " Left/right panels stream live translations, enabling students to switch languages dynamically.")
]

for title, desc in features_d:
    p = tf_d.add_paragraph()
    p.space_after = Pt(12)
    run_title = p.add_run()
    run_title.text = "• " + title
    run_title.font.bold = True
    run_title.font.size = Pt(14)
    run_title.font.color.rgb = WHITE
    run_desc = p.add_run()
    run_desc.text = desc
    run_desc.font.size = Pt(13)
    run_desc.font.color.rgb = EMERALD

# Right Image Mockup
if os.path.exists(DASHBOARD_DEST):
    slide4.shapes.add_picture(DASHBOARD_DEST, Inches(6.0), Inches(1.6), Inches(6.5), Inches(4.85))
else:
    # Fallback shape if file missing
    fallback = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.0), Inches(1.6), Inches(6.5), Inches(4.85))
    fallback.fill.solid()
    fallback.fill.fore_color.rgb = DARK_BLUE_LIGHT
    fallback.line.color.rgb = EMERALD

# ==========================================
# SLIDE 5: SYSTEM ARCHITECTURE (LIGHT MODE)
# ==========================================
slide5 = prs.slides.add_slide(blank_layout)
add_full_background(slide5, OFF_WHITE)
add_header(slide5, "System Architecture & Processing Pipeline")

# Let's draw 3 architectural block columns:
# 1. Inputs (Sources)
# 2. AI & Processing Engine (FastAPI backend)
# 3. Output Delivery Layer (User Interfaces)

# 1. Input Box
add_card(slide5, Inches(0.8), Inches(1.8), Inches(3.2), Inches(4.5), WHITE, border_color=DARK_BLUE)
in_tb = slide5.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(2.8), Inches(4.1))
tf_in = in_tb.text_frame
tf_in.word_wrap = True
pin = tf_in.paragraphs[0]
pin.text = "1. INGESTION LAYER"
pin.font.bold = True
pin.font.size = Pt(16)
pin.font.color.rgb = DARK_BLUE
pin.space_after = Pt(15)

in_pts = ["• Live Classroom Audio Feed", "• PDFs & Textbook Scans", "• Youtube Lecture URLs", "• WebRTC Voice Commands", "• Peer Tutors Inputs"]
for pt in in_pts:
    p = tf_in.add_paragraph()
    p.text = pt
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_DARK
    p.space_after = Pt(10)

# Connector Arrow 1
arrow1 = slide5.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.2), Inches(3.8), Inches(0.7), Inches(0.4))
arrow1.fill.solid()
arrow1.fill.fore_color.rgb = EMERALD
arrow1.line.color.rgb = EMERALD

# 2. Backend Box
add_card(slide5, Inches(5.1), Inches(1.8), Inches(3.2), Inches(4.5), DARK_BLUE)
mid_tb = slide5.shapes.add_textbox(Inches(5.3), Inches(2.0), Inches(2.8), Inches(4.1))
tf_mid = mid_tb.text_frame
tf_mid.word_wrap = True
pmid = tf_mid.paragraphs[0]
pmid.text = "2. CORE AI PIPELINE"
pmid.font.bold = True
pmid.font.size = Pt(16)
pmid.font.color.rgb = EMERALD
pmid.space_after = Pt(15)

mid_pts = ["• Web Speech API (SpeechRecognition)", "• Google Gemini 2.5 Flash SDK", "• Dynamic Node Expansion logic", "• Web Speech API (SpeechSynthesis TTS)", "• Express Node.js Backend Server"]
for pt in mid_pts:
    p = tf_mid.add_paragraph()
    p.text = pt
    p.font.size = Pt(13)
    p.font.color.rgb = WHITE
    p.space_after = Pt(10)

# Connector Arrow 2
arrow2 = slide5.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(8.5), Inches(3.8), Inches(0.7), Inches(0.4))
arrow2.fill.solid()
arrow2.fill.fore_color.rgb = EMERALD
arrow2.line.color.rgb = EMERALD

# 3. Output Box
add_card(slide5, Inches(9.4), Inches(1.8), Inches(3.1), Inches(4.5), WHITE, border_color=EMERALD)
out_tb = slide5.shapes.add_textbox(Inches(9.6), Inches(2.0), Inches(2.7), Inches(4.1))
tf_out = out_tb.text_frame
tf_out.word_wrap = True
pout = tf_out.paragraphs[0]
pout.text = "3. PRESENTATION LAYER"
pout.font.bold = True
pout.font.size = Pt(16)
pout.font.color.rgb = DARK_BLUE
pout.space_after = Pt(15)

out_pts = ["• Interactive Canvas Nodes", "• Adaptive Study Cards", "• Real-Time Audio Player", "• Screen-Reader DOM Layout", "• Peer-to-Peer Rewards Hub"]
for pt in out_pts:
    p = tf_out.add_paragraph()
    p.text = pt
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_DARK
    p.space_after = Pt(10)


# ==========================================
# SLIDE 6: TECHNOLOGY STACK (LIGHT MODE)
# ==========================================
slide6 = prs.slides.add_slide(blank_layout)
add_full_background(slide6, OFF_WHITE)
add_header(slide6, "The Scalable Technology Stack")

stacks = [
    ("FRONTEND (REACT)", "Highly accessible, interactive user interface.", 
     ["• React.js & TypeScript", "• Tailwind CSS (Glassmorphism)", "• HTML5 Canvas (VisualMind Graph)", "• Web Speech API (STT Transcription)", "• Web Speech API (TTS Vocalizer)"], 
     Inches(0.8), Inches(1.8)),
     
    ("BACKEND & AI", "Express middleware and Gemini LLM integrations.", 
     ["• Node.js & Express.js", "• Google Gemini 2.5 Flash SDK", "• Multer File Ingestor", "• Dynamic Quiz Generator", "• Dynamic Translation endpoints"], 
     Inches(4.9), Inches(1.8)),
     
    ("UTILITIES & DEPLOY", "Staging, documentation, and cloud execution.", 
     ["• reportlab (PDF Compiler)", "• python-pptx (Deck Generator)", "• Git & GitHub Version Control", "• Render Cloud Web Services", "• Unified MERN static delivery"], 
     Inches(9.0), Inches(1.8))
]

for title, desc, tools, left, top in stacks:
    add_card(slide6, left, top, Inches(3.6), Inches(4.8), WHITE, border_color=EMERALD)
    tb = slide6.shapes.add_textbox(left + Inches(0.15), top + Inches(0.2), Inches(3.3), Inches(4.4))
    tf_st = tb.text_frame
    tf_st.word_wrap = True
    
    p1 = tf_st.paragraphs[0]
    p1.text = title
    p1.font.name = "Arial"
    p1.font.bold = True
    p1.font.size = Pt(16)
    p1.font.color.rgb = DARK_BLUE
    p1.space_after = Pt(6)
    
    p2 = tf_st.add_paragraph()
    p2.text = desc
    p2.font.name = "Arial"
    p2.font.italic = True
    p2.font.size = Pt(11)
    p2.font.color.rgb = TEXT_MUTED
    p2.space_after = Pt(14)
    
    for tool in tools:
        p = tf_st.add_paragraph()
        p.text = tool
        p.font.name = "Arial"
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(8)


# ==========================================
# SLIDE 7: EDUACCESS & IMPACT (LIGHT MODE)
# ==========================================
slide7 = prs.slides.add_slide(blank_layout)
add_full_background(slide7, OFF_WHITE)
add_header(slide7, "EduAccess & Social Impact")

# Left Image Mockup
if os.path.exists(ACCESSIBILITY_DEST):
    slide7.shapes.add_picture(ACCESSIBILITY_DEST, Inches(0.8), Inches(1.6), Inches(5.6), Inches(4.8))
else:
    fallback = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.6), Inches(5.6), Inches(4.8))
    fallback.fill.solid()
    fallback.fill.fore_color.rgb = DARK_BLUE_LIGHT
    fallback.line.color.rgb = EMERALD

# Right Text Card
add_card(slide7, Inches(6.8), Inches(1.6), Inches(5.7), Inches(4.8), WHITE, border_color=DARK_BLUE)
imp_tb = slide7.shapes.add_textbox(Inches(7.1), Inches(1.8), Inches(5.1), Inches(4.4))
tf_imp = imp_tb.text_frame
tf_imp.word_wrap = True

pimp1 = tf_imp.paragraphs[0]
pimp1.text = "Designing for Universal Inclusion:"
pimp1.font.name = "Arial"
pimp1.font.size = Pt(20)
pimp1.font.bold = True
pimp1.font.color.rgb = DARK_BLUE
pimp1.space_after = Pt(16)

impact_pts = [
    ("Voice Control Integration:", " Screen-free interaction models allowing visually impaired students to dictate notes, navigate screens, and search topic nodes using spoken queries."),
    ("Neurodiverse Mode:", " Auto-simplifies dense jargon into structured, high-contrast summaries with minimized distraction elements."),
    ("Expected Impact:", " Bridges language-medium gaps for up to 10M+ Indian students annually; reduces prep time for exam prep by 40%; provides a democratic level playing field for diverse abilities.")
]

for title, desc in impact_pts:
    p = tf_imp.add_paragraph()
    p.space_after = Pt(10)
    run_title = p.add_run()
    run_title.text = "• " + title
    run_title.font.bold = True
    run_title.font.size = Pt(13)
    run_title.font.color.rgb = EMERALD_MUTED
    run_desc = p.add_run()
    run_desc.text = desc
    run_desc.font.size = Pt(13)
    run_desc.font.color.rgb = TEXT_DARK


# ==========================================
# SLIDE 8: ROADMAP & FUTURE WORK (DARK MODE)
# ==========================================
slide8 = prs.slides.add_slide(blank_layout)
add_full_background(slide8, DARK_BLUE)
add_header(slide8, "Development Roadmap & Milestones", dark_mode=True)

# Define 4 horizontal cards along a timeline
timeline_phases = [
    ("ROUND 1", "June 17 - 18", "• Problem Statement identification\n• Idea formulation & PPTX deck\n• Concept validation\n• Proposal Submission", Inches(0.8)),
    ("ROUND 2", "June 20 - 22", "• Express/Node backend config\n• Whisper & Gemini integrations\n• Interactive canvas UI prototype\n• Source code & Video submission", Inches(3.8)),
    ("PHASE 3", "June 23 - 26", "• Pilot testing & bug fixes\n• Real-time socket connections\n• Demo polishing & presentation compilation", Inches(6.8)),
    ("GRAND FINALE", "June 27 - 28", "• Deployment of final MERN app\n• Pitch presentation to expert jury\n• Q&A & Design defense", Inches(9.8))
]

# Draw horizontal connector line
conn_line = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(3.5), Inches(10.0), Inches(0.08))
conn_line.fill.solid()
conn_line.fill.fore_color.rgb = EMERALD_MUTED
conn_line.line.color.rgb = EMERALD_MUTED

for title, time_lbl, details, left in timeline_phases:
    # Card background
    add_card(slide8, left, Inches(1.8), Inches(2.7), Inches(4.5), DARK_BLUE_LIGHT, border_color=EMERALD)
    
    # Text
    tb = slide8.shapes.add_textbox(left + Inches(0.1), Inches(2.0), Inches(2.5), Inches(4.1))
    tf_rd = tb.text_frame
    tf_rd.word_wrap = True
    
    p1 = tf_rd.paragraphs[0]
    p1.text = title
    p1.font.bold = True
    p1.font.size = Pt(18)
    p1.font.color.rgb = EMERALD
    p1.space_after = Pt(4)
    
    p2 = tf_rd.add_paragraph()
    p2.text = time_lbl
    p2.font.italic = True
    p2.font.size = Pt(12)
    p2.font.color.rgb = WHITE
    p2.space_after = Pt(15)
    
    p3 = tf_rd.add_paragraph()
    p3.text = details
    p3.font.size = Pt(12)
    p3.font.color.rgb = WHITE

# Add closing footer note
footer_box = slide8.shapes.add_textbox(Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.6))
tf_ft = footer_box.text_frame
tf_ft.word_wrap = True
pft = tf_ft.paragraphs[0]
pft.text = "Bharat Academix CodeQuest 2026  |  Empowering academic excellence through modern, accessible AI systems."
pft.alignment = PP_ALIGN.CENTER
pft.font.size = Pt(11)
pft.font.color.rgb = TEXT_MUTED

# Save the presentation with fallback if file is locked
output_base = "AcademixIQ_Pitch_Deck"
saved = False
for suffix in ["", "_v2", "_v3", "_v4", "_v5"]:
    current_name = f"{output_base}{suffix}.pptx"
    try:
        prs.save(current_name)
        print(f"Presentation saved successfully as {current_name}!")
        saved = True
        break
    except PermissionError:
        print(f"Warning: {current_name} is locked/open. Trying next slot...")

if not saved:
    import time
    fallback_name = f"{output_base}_{int(time.time())}.pptx"
    prs.save(fallback_name)
    print(f"Presentation saved successfully as fallback {fallback_name}!")
